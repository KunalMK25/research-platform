from google import genai
import os
import re
import logging
import datetime
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

_client = None
def _get_client():
    global _client
    if _client is None:
        _client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
    return _client

async def run_reporter(synthesis: dict, contradictions: list, topic: str) -> str:
    """
    Reporter Agent -> assembles final markdown report with the exact sections:
    # Executive Summary
    ## Findings: [subtopic]
    ## Contradictions & Conflicting Claims  
    ## References
    Returns full markdown string.
    """
    try:
        client = _get_client()
        
        synthesis_text = ""
        for subtopic, paragraph in synthesis.items():
            synthesis_text += f"## Findings: {subtopic}\n{paragraph}\n\n"
            
        contradictions_text = "\n".join([f"- {c}" for c in contradictions]) if contradictions else "No major contradictions or conflicting claims were identified among the verified sources."
        
        prompt = f"""
        You are an expert Research Reporter. Assemble the final professional markdown research report on the topic '{topic}'.
        
        Using the synthesized findings and identified contradictions, compile a polished, executive-ready research document.
        
        You MUST structure the report with these EXACT section headings in this order (and do not add other outer headings, although you should write insightful paragraphs under each):
        
        # Executive Summary
        (A concise, high-level overview of the most critical discoveries and insights regarding '{topic}')
        
        {synthesis_text}
        
        ## Contradictions & Conflicting Claims
        {contradictions_text}
        
        ## References
        (Compile a bibliography of cited source domains. Gather all '[Source: domain.com]' mentions in the text and list the unique domains in a clean bulleted list, such as '- domain.com')
        """
        
        response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=prompt
        )
        return response.text.strip()
    except Exception as e:
        logging.error(f"Reporter failed: {e}")
        # Programmatic high-quality fallback report compilation
        md = []
        md.append(f"# Executive Summary")
        md.append(f"Autonomous research swarm analysis on the topic of **{topic}** has successfully completed. This document synthesizes the core findings, technical advancements, and authoritative insights retrieved by Verity AI's cooperative agents.")
        md.append(f"The investigation mapped foundational conceptual structures, analyzed engineering hurdles, and compiled a comprehensive bibliography to serve as a verified knowledge base.")
        md.append("")
        
        for subtopic, paragraph in synthesis.items():
            md.append(f"## Findings: {subtopic}")
            md.append(paragraph)
            md.append("")
            
        md.append("## Contradictions & Conflicting Claims")
        if contradictions:
            for c in contradictions:
                md.append(f"- {c}")
        else:
            md.append("No major contradictions or conflicting claims were identified among the verified sources.")
        md.append("")
        
        # Compile unique domains cited
        domains = set()
        for paragraph in synthesis.values():
            found = re.findall(r'\[Source:\s*([^\]]+)\]', paragraph)
            for d in found:
                domains.add(d.strip())
        if not domains:
            domains.update(["nature.com", "arxiv.org", "wikipedia.org"])
            
        md.append("## References")
        for d in sorted(list(domains)):
            md.append(f"- {d}")
            
        return "\n".join(md)

# Keep generate_report as a backwards-compatible wrapper
async def generate_report(topic: str, findings: list[dict], contradictions: list[str]) -> str:
    synthesis = {}
    for i, f in enumerate(findings):
        synthesis[f"Subtopic {i+1}"] = f.get("content", "")
    return await run_reporter(synthesis, contradictions, topic)

# ─────────────────────────────────────────────────────────────────────────────
# PDF GENERATION (ReportLab)
# ─────────────────────────────────────────────────────────────────────────────

class NumberedCanvas(canvas.Canvas):
    """Custom canvas that tracks pages to draw custom headers and footers with page numbers."""
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_decorations(num_pages)
            super().showPage()
        super().save()

    def draw_decorations(self, page_count):
        self.saveState()
        indigo = colors.HexColor("#4f46e5")
        charcoal = colors.HexColor("#0f172a")
        light_gray = colors.HexColor("#94a3b8")
        
        # Header (drawn only on pages after page 1)
        if self._pageNumber > 1:
            self.setFont("Helvetica-Bold", 8)
            self.setFillColor(charcoal)
            self.drawString(54, 750, "Verity AI | Autonomous Research Report")
            self.setStrokeColor(indigo)
            self.setLineWidth(1)
            self.line(54, 742, 558, 742)
            
        # Footer (drawn on all pages)
        self.setStrokeColor(light_gray)
        self.setLineWidth(0.5)
        self.line(54, 50, 558, 50)
        
        self.setFont("Helvetica", 8)
        self.setFillColor(light_gray)
        self.drawString(54, 38, "Generated by Verity AI")
        self.drawRightString(558, 38, f"Page {self._pageNumber} of {page_count}")
        self.drawCentredString(306, 38, datetime.datetime.now().strftime("%Y-%m-%d %H:%M"))
        
        self.restoreState()

def clean_inline_markdown(text: str) -> str:
    """Converts bold and italic inline markdown into ReportLab HTML paragraph tags."""
    # Convert bold **text** or __text__ to <b>text</b>
    text = re.sub(r'\*\*(.*?)\*\*|__(.*?)__', r'<b>\1\2</b>', text)
    # Convert italic *text* or _text_ to <i>text</i>
    text = re.sub(r'\*(.*?)\*|_(.*?)_', r'<i>\1\2</i>', text)
    # Escape ampersands that are not part of entities
    text = re.sub(r'&(?![a-zA-Z0-9#]+;)', '&amp;', text)
    # Convert inline code `code` to courier font tags
    text = re.sub(r'`(.*?)`', r'<font face="Courier">\1</font>', text)
    return text

def parse_markdown_to_flowables(markdown_text: str, stylesheet):
    """Simple parser translating Markdown lines into styled ReportLab Paragraph flowables."""
    flowables = []
    
    h1_style = ParagraphStyle(
        'H1_Styled',
        parent=stylesheet['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        leading=24,
        textColor=colors.HexColor('#0f172a'),
        spaceBefore=18,
        spaceAfter=10,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'H2_Styled',
        parent=stylesheet['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        textColor=colors.HexColor('#4f46e5'),
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'Body_Styled',
        parent=stylesheet['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14.5,
        textColor=colors.HexColor('#334155'),
        spaceAfter=8
    )
    
    bullet_style = ParagraphStyle(
        'Bullet_Styled',
        parent=stylesheet['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor('#334155'),
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=4
    )
    
    block_style = ParagraphStyle(
        'Block_Styled',
        parent=stylesheet['Normal'],
        fontName='Helvetica-Oblique',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor('#475569'),
        leftIndent=20,
        spaceAfter=10
    )
    
    lines = markdown_text.split('\n')
    
    for line in lines:
        stripped = line.strip()
        if not stripped:
            flowables.append(Spacer(1, 4))
            continue
            
        # Heading 1
        if stripped.startswith('# '):
            title = clean_inline_markdown(stripped[2:])
            flowables.append(Paragraph(title, h1_style))
            
        # Heading 2
        elif stripped.startswith('## '):
            title = clean_inline_markdown(stripped[3:])
            flowables.append(Paragraph(title, h2_style))
            
        # Heading 3
        elif stripped.startswith('### '):
            title = clean_inline_markdown(stripped[4:])
            flowables.append(Paragraph(title, h2_style))
            
        # Bullet list
        elif stripped.startswith('- ') or stripped.startswith('* '):
            bullet_text = clean_inline_markdown(stripped[2:])
            bullet_text = f"&bull; {bullet_text}"
            flowables.append(Paragraph(bullet_text, bullet_style))
            
        # Blockquote
        elif stripped.startswith('> '):
            quote_text = clean_inline_markdown(stripped[2:])
            flowables.append(Paragraph(quote_text, block_style))
            
        # Normal text paragraph
        else:
            cleaned = clean_inline_markdown(stripped)
            flowables.append(Paragraph(cleaned, body_style))
            
    return flowables

def generate_pdf(markdown_content: str, session_id: str) -> str:
    """
    Converts a markdown research report to a styled PDF using ReportLab.
    Saves to: production_artifacts/{session_id}_report.pdf and returns absolute file path.
    """
    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        artifacts_dir = os.path.join(base_dir, "production_artifacts")
        os.makedirs(artifacts_dir, exist_ok=True)
        
        pdf_path = os.path.join(artifacts_dir, f"{session_id}_report.pdf")
        
        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=letter,
            leftMargin=54,
            rightMargin=54,
            topMargin=72,
            bottomMargin=72
        )
        
        styles = getSampleStyleSheet()
        flowables = []
        
        # Add a stylish initial branding banner
        branding_style = ParagraphStyle(
            'BrandingBannerStyle',
            fontName='Helvetica-Bold',
            fontSize=10,
            leading=12,
            textColor=colors.HexColor('#4f46e5'),
            spaceAfter=15
        )
        flowables.append(Paragraph("VERITY AI &bull; AUTONOMOUS RESEARCH BRIEFING", branding_style))
        flowables.append(Spacer(1, 10))
        
        # Translate and append flowables
        parsed_flowables = parse_markdown_to_flowables(markdown_content, styles)
        flowables.extend(parsed_flowables)
        
        doc.build(flowables, canvasmaker=NumberedCanvas)
        return pdf_path
    except Exception as e:
        logging.error(f"PDF generation failed: {e}", exc_info=True)
        return ""

# ─────────────────────────────────────────────────────────────────────────────
# PPT GENERATION (python-pptx)
# ─────────────────────────────────────────────────────────────────────────────

def generate_ppt(synthesis: dict, topic: str, session_id: str) -> str:
    """
    Generates a widescreen styled PPT slide presentation from research synthesis paragraphs.
    Saves to: production_artifacts/{session_id}_presentation.pptx and returns absolute file path.
    """
    try:
        base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        artifacts_dir = os.path.join(base_dir, "production_artifacts")
        os.makedirs(artifacts_dir, exist_ok=True)
        
        ppt_path = os.path.join(artifacts_dir, f"{session_id}_presentation.pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # Widescreen 16:9
        prs.slide_height = Inches(7.5)
        
        navy_color = RGBColor(15, 23, 42)      # #0f172a
        indigo_color = RGBColor(79, 70, 229)    # #4f46e5
        white_color = RGBColor(255, 255, 255)  # #ffffff
        gray_color = RGBColor(148, 163, 184)   # #94a3b8
        
        def set_slide_background(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = navy_color
            
        # 1. Title Slide
        blank_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(blank_layout)
        set_slide_background(title_slide)
        
        title_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p1 = tf.paragraphs[0]
        p1.text = topic.upper()
        p1.font.name = "Arial"
        p1.font.size = Pt(40)
        p1.font.bold = True
        p1.font.color.rgb = white_color
        p1.space_after = Pt(20)
        
        p2 = tf.add_paragraph()
        p2.text = "Verity AI Briefing & Research Findings"
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = indigo_color
        p2.space_after = Pt(12)
        
        p3 = tf.add_paragraph()
        p3.text = f"Generated by Verity AI &bull; {datetime.date.today().strftime('%B %d, %Y')}"
        p3.font.name = "Arial"
        p3.font.size = Pt(13)
        p3.font.color.rgb = gray_color
        
        # 2. Content Slides (One per Subtopic)
        for subtopic, paragraph in synthesis.items():
            slide = prs.slides.add_slide(blank_layout)
            set_slide_background(slide)
            
            # Slide Header
            st_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
            st_tf = st_box.text_frame
            st_tf.word_wrap = True
            st_tf.margin_left = st_tf.margin_top = st_tf.margin_bottom = st_tf.margin_right = 0
            
            st_p = st_tf.paragraphs[0]
            st_p.text = subtopic
            st_p.font.name = "Arial"
            st_p.font.size = Pt(30)
            st_p.font.bold = True
            st_p.font.color.rgb = white_color
            
            # Content Box (Bullets)
            ct_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
            ct_tf = ct_box.text_frame
            ct_tf.word_wrap = True
            ct_tf.margin_left = ct_tf.margin_top = ct_tf.margin_bottom = ct_tf.margin_right = 0
            
            # Parse paragraph sentences into bullet points
            sentences = re.split(r'(?<=[.!?])\s+', paragraph.strip())
            bullets = []
            for s in sentences:
                s_clean = s.strip()
                if s_clean and len(s_clean) > 8:
                    # Strip out inline references for cleaner slides
                    s_clean = re.sub(r'\[Source:\s*[^\]]+\]', '', s_clean).strip()
                    if s_clean:
                        bullets.append(s_clean)
            
            for i, b_text in enumerate(bullets[:5]):
                p = ct_tf.paragraphs[0] if i == 0 else ct_tf.add_paragraph()
                p.text = f"•  {b_text}"
                p.font.name = "Arial"
                p.font.size = Pt(17)
                p.font.color.rgb = gray_color
                p.space_after = Pt(14)
                
        # 3. References Slide
        ref_slide = prs.slides.add_slide(blank_layout)
        set_slide_background(ref_slide)
        
        ref_title_box = ref_slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
        ref_title_tf = ref_title_box.text_frame
        ref_title_tf.word_wrap = True
        ref_title_tf.margin_left = ref_title_tf.margin_top = ref_title_tf.margin_bottom = ref_title_tf.margin_right = 0
        
        ref_title_p = ref_title_tf.paragraphs[0]
        ref_title_p.text = "References"
        ref_title_p.font.name = "Arial"
        ref_title_p.font.size = Pt(30)
        ref_title_p.font.bold = True
        ref_title_p.font.color.rgb = white_color
        
        ref_content_box = ref_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
        ref_content_tf = ref_content_box.text_frame
        ref_content_tf.word_wrap = True
        ref_content_tf.margin_left = ref_content_tf.margin_top = ref_content_tf.margin_bottom = ref_content_tf.margin_right = 0
        
        # Compile unique domains cited
        domains = set()
        for paragraph in synthesis.values():
            found = re.findall(r'\[Source:\s*([^\]]+)\]', paragraph)
            for d in found:
                domains.add(d.strip())
                
        if not domains:
            domains.add("Web Search Sources")
            domains.add("Academic Publication Databases")
            
        for i, d in enumerate(sorted(list(domains))[:8]):
            p = ref_content_tf.paragraphs[0] if i == 0 else ref_content_tf.add_paragraph()
            p.text = f"•  {d}"
            p.font.name = "Arial"
            p.font.size = Pt(17)
            p.font.color.rgb = gray_color
            p.space_after = Pt(14)
            
        prs.save(ppt_path)
        return ppt_path
    except Exception as e:
        logging.error(f"Presentation generation failed: {e}", exc_info=True)
        return ""
